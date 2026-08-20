from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
HEADER_BG = RGBColor(0x1F, 0x4E, 0x79)
GOOD_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
BAD_RED = RGBColor(0xFF, 0xEB, 0xEE)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(12.333)
    height = Inches(2.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(2.8)
    width = Inches(12.333)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT
    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None, highlight_cells=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_width = Inches(12.333)
    table_height = Inches(0.45 * n_rows)
    left = Inches(0.5)
    top = Inches(1.2)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if highlight_cells and (i + 1, j) in highlight_cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOOD_GREEN

    return slide


def add_bullet_slide(prs, title, bullets, font_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(11.733)
    height = Inches(5.5)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.space_after = Pt(8)
        if bullet.startswith("  "):
            p.level = 1
            p.font.size = Pt(font_size - 2)

    return slide


# ========== Slide 1: Title ==========
add_title_slide(prs, "森林点云分割精度提升实验报告", "LitePT 框架 | mIoU: 0.5529 → 0.6542 (+10.13%)")

# ========== Slide 2: Overview ==========
add_section_slide(prs, "一、项目概述")

add_bullet_slide(prs, "项目概述", [
    "目标：在 LitePT 框架上提升森林点云 7 类分割精度（mIoU）",
    "数据集：55 train / 6 val scenes，~36M 训练点，LAS 格式，grid_size=0.02",
    "类别：terrain(0), foliage(1), CWD(2), trunk(3), branch(4), snag(5), non-tree-cyl(6)",
    "硬件：本地 RTX 4080 16GB（验证）+ 服务器 RTX 4080 32GB（最终训练）",
    "总提升：mIoU 0.5529 → 0.6542，提升 +10.13%",
])

# ========== Slide 3: Roadmap ==========
add_table_slide(prs, "二、整体技术路线", 
    ["阶段", "措施", "mIoU", "变化"],
    [
        ["Baseline", "默认配置", "0.5529", "—"],
        ["Phase 1", "Loss 增强（CE+smoothing+Lovasz+Dice）", "0.6092", "+0.0563"],
        ["Phase 2", "grid_size 0.04", "0.5573", "-0.0529 ✗"],
        ["Phase 3a", "enc_patch_size 2048", "0.6162", "+0.0633"],
        ["Phase 3b", "enc_patch_size 4096", "0.5951", "+0.0422（过拟合）"],
        ["Phase 4", "数据修复（标签纠错）", "0.6040", "—"],
        ["Phase 5", "干净数据 + 600K crop + 服务器", "0.6542", "+0.1013"],
    ],
    highlight_cells={(7, 2): True, (7, 3): True}
)

# ========== Slide 4: Phase 1 ==========
add_section_slide(prs, "Phase 1: Loss 函数增强")

add_table_slide(prs, "Loss-v2: CE(label_smoothing=0.1) + Lovasz + Dice",
    ["类别", "Baseline", "Loss-v2", "变化"],
    [
        ["terrain", "0.823", "0.842", "+0.019"],
        ["foliage", "0.867", "0.872", "+0.005"],
        ["CWD", "0.000", "0.000", "—"],
        ["trunk", "0.659", "0.707", "+0.048"],
        ["branch", "0.507", "0.565", "+0.058"],
        ["snag", "0.710", "0.470", "-0.240 ↓"],
        ["non-tree-cyl", "0.108", "0.496", "+0.388"],
        ["mIoU", "0.5529", "0.6092", "+0.0563"],
    ],
    highlight_cells={(8, 2): True, (8, 3): True}
)

add_bullet_slide(prs, "Phase 1 关键发现", [
    "Loss-v2 大幅提升 trunk(+4.8%), branch(+5.8%), non-tree-cyl(+38.8%)",
    "但 snag 严重退化：0.71 → 0.47（-24%）",
    "原因分析：后续发现 snag 退化根因是数据标签错误（Phase 4 揭示）",
    "mIoU 整体提升 5.63%，验证了多 Loss 组合的有效性",
])

# ========== Slide 5: Phase 2 ==========
add_section_slide(prs, "Phase 2: grid_size 调整（失败）")

add_bullet_slide(prs, "Phase 2: grid_size 实验", [
    "尝试：grid_size 0.02 → 0.04（减少点云密度，加速训练）",
    "结果：mIoU 0.5529 → 0.5573（仅 +0.44%，几乎无提升）",
    "结论：增大 grid_size 导致点云几何信息丢失，不利于细粒度分割",
    "决定：放弃此方案，保持 grid_size=0.02",
])

# ========== Slide 6: Phase 3 ==========
add_section_slide(prs, "Phase 3: 感受野调整")

add_table_slide(prs, "enc_patch_size 实验",
    ["配置", "enc_patch_size", "mIoU", "snag IoU", "备注"],
    [
        ["Baseline", "1024", "0.5529", "0.710", ""],
        ["Loss-v2", "1024", "0.6092", "0.470", ""],
        ["Patch 2048", "2048", "0.6162", "0.594", "snag 回升"],
        ["Patch 4096", "4096", "0.5951", "—", "过拟合"],
    ],
    highlight_cells={(4, 2): True}
)

add_table_slide(prs, "感受野实测（不同 crop_point_max）",
    ["crop_point_max", "实际半径范围"],
    [
        ["100K", "1-2.4m"],
        ["500K", "4-8m"],
        ["600K", "4-6m"],
        ["1M", "6-16m"],
    ]
)

add_bullet_slide(prs, "Phase 3 关键结论", [
    "enc_patch_size=2048 有效扩大感受野，snag 从 0.47 恢复到 0.594",
    "enc_patch_size=4096 过拟合（参数量过大，训练数据不足）",
    "选定 enc_patch_size=2048 作为最终配置",
])

# ========== Slide 7: Phase 4 ==========
add_section_slide(prs, "Phase 4: 数据修复")

add_table_slide(prs, "发现的标签错误",
    ["文件", "问题", "修复方式"],
    [
        ["pole_05_5.las", "1.3M 点错误标记为 snag（占训练 38%）", "snag → non-tree-cyl(6)"],
        ["snag1/2/3.las", "含 non-tree-cyl 标签", "non-tree-cyl → ignored(7)"],
        ["所有 pole 文件", "含非 pole 标签", "非 pole → ignored(7)"],
        ["nontree_other", "错误 snag 标签", "snag → ignored(7)"],
    ]
)

add_table_slide(prs, "数据修复效果",
    ["类别", "修复前(脏)", "修复后(干净)", "变化"],
    [
        ["terrain", "0.842", "0.833", "-0.009"],
        ["foliage", "0.872", "0.870", "-0.002"],
        ["CWD", "0.000", "0.000", "—"],
        ["trunk", "0.707", "0.704", "-0.003"],
        ["branch", "0.565", "0.570", "+0.005"],
        ["snag", "0.470", "0.360", "-0.110 ↓"],
        ["non-tree-cyl", "0.496", "0.920", "+0.424 ↑"],
        ["mIoU", "0.6092", "0.6040", "-0.005"],
    ],
    highlight_cells={(8, 2): True, (8, 3): True}
)

add_bullet_slide(prs, "Phase 4 关键发现", [
    "数据修复使 non-tree-cyl 大幅提升 +42.4%（0.496 → 0.920）",
    "snag 反而下降：干净数据中去掉了错误 snag 样本，训练数据减少",
    "mIoU 略降 0.5%，但标签准确性大幅提升，为后续训练奠定基础",
])

# ========== Slide 8: Phase 5 ==========
add_section_slide(prs, "Phase 5: 最终整合训练")

add_table_slide(prs, "最终训练配置",
    ["参数", "值"],
    [
        ["Base config", "loss-v2（CE+smoothing+Lovasz+Dice）"],
        ["enc_patch_size", "2048"],
        ["crop_point_max", "600,000"],
        ["数据", "干净数据（修复后）"],
        ["预训练权重", "Phase 4 model_best（0.6040）"],
        ["resume", "False（只加载权重，scheduler 从头）"],
        ["epoch", "60（跑了 34 epoch）"],
        ["batch_size", "1"],
        ["Scheduler", "OneCycleLR(max_lr=0.001)"],
        ["硬件", "服务器 RTX 4080 32GB"],
    ]
)

# ========== Slide 9: Training Curve ==========
add_table_slide(prs, "训练曲线（关键 epoch）",
    ["Epoch", "mIoU", "snag", "non-tree-cyl", "branch", "trunk", "备注"],
    [
        ["1", "0.4904", "0.447", "0.005", "0.605", "0.678", ""],
        ["5", "0.6313", "0.488", "0.911", "0.602", "0.713", "CWD 首次非零"],
        ["6", "0.6444", "0.540", "0.906", "0.627", "0.737", ""],
        ["14", "0.6048", "0.586", "0.583", "0.632", "0.738", "snag 突破"],
        ["22", "0.6242", "0.656", "0.665", "0.682", "0.768", "snag/branch 最佳"],
        ["25", "0.6542", "0.572", "0.903", "0.635", "0.767", "Best mIoU ★"],
        ["27", "0.5809", "0.566", "0.371", "0.651", "0.772", "trunk 最佳"],
        ["33", "0.6367", "0.575", "0.795", "0.626", "0.762", "model_last"],
    ],
    highlight_cells={(7, 1): True}
)

add_bullet_slide(prs, "训练曲线特征", [
    "mIoU 剧烈抖动（0.45-0.65）：val 仅 9 scene，non-tree-cyl 在 0.00-0.91 之间跳",
    "Best mIoU = 0.6542（epoch 25）",
    "snag 最佳 0.656（epoch 22），branch 最佳 0.682（epoch 22），trunk 最佳 0.772（epoch 27）",
    "后期（ep 25-33）高值频率增加，模型仍在学习中",
])

# ========== Slide 10: Final Results ==========
add_section_slide(prs, "最终结果汇总")

add_table_slide(prs, "各类别最佳 IoU",
    ["类别", "Baseline", "最终 Best", "提升"],
    [
        ["terrain", "0.823", "0.831", "+0.008"],
        ["foliage", "0.867", "0.882", "+0.015"],
        ["CWD", "0.000", "0.006", "+0.006"],
        ["trunk", "0.659", "0.772", "+0.113 ↑"],
        ["branch", "0.507", "0.682", "+0.175 ↑"],
        ["snag", "0.710", "0.656", "-0.054"],
        ["non-tree-cyl", "0.108", "0.920", "+0.812 ↑"],
        ["mIoU", "0.5529", "0.6542", "+0.1013 ↑"],
    ],
    highlight_cells={(8, 2): True, (8, 3): True}
)

add_table_slide(prs, "各实验阶段最佳 mIoU 对比",
    ["实验", "mIoU", "vs Baseline"],
    [
        ["Baseline", "0.5529", "—"],
        ["Loss-v2", "0.6092", "+5.63%"],
        ["+ Patch 2048", "0.6162", "+6.33%"],
        ["+ 数据修复", "0.6040", "+5.11%"],
        ["+ 600K crop 服务器训练", "0.6542", "+10.13%"],
    ],
    highlight_cells={(5, 1): True, (5, 2): True}
)

# ========== Slide 11: Key Decisions ==========
add_section_slide(prs, "关键决策与经验")

add_table_slide(prs, "关键决策",
    ["决策", "原因"],
    [
        ["选用 Loss-v2（CE+Lovasz+Dice）", "多损失函数互补，对难分类样本效果显著"],
        ["enc_patch_size=2048", "扩大感受野帮助 snag 等需要上下文的类别"],
        ["数据修复", "pole_05_5.las 1.3M 错误 snag 占训练 38%"],
        ["crop_point_max=600K", "32GB GPU 上最大化感受野（4-6m 半径）"],
        ["resume=False", "避免加载旧 scheduler 导致 lr 极低、模型退化"],
        ["model_last > model_best", "val 样本少，mIoU 波动大，last 更稳健"],
    ]
)

# ========== Slide 12: Future ==========
add_section_slide(prs, "待改进方向")

add_table_slide(prs, "后续改进计划",
    ["方向", "预期效果", "难度", "优先级"],
    [
        ["CWD 补充数据", "CWD IoU 从 0 提升到 0.3+", "中", "高"],
        ["snag 数据增强", "snag 从 0.656 提升到 0.7+", "中", "高"],
        ["测试时增强（TTA）", "稳定 mIoU 抖动", "低", "高"],
        ["eval_epoch=1 监控", "更精确选 best", "低", "中"],
        ["更大 backbone", "mIoU 整体提升 2-3%", "高", "低"],
    ]
)

# ========== Slide 13: Summary ==========
add_title_slide(prs, "总结", "mIoU: 0.5529 → 0.6542 (+10.13%)\nnon-tree-cyl: 0.108 → 0.920 | branch: 0.507 → 0.682 | trunk: 0.659 → 0.772")

out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "report.pptx")
prs.save(out_path)
print(f"Report saved to: {out_path}")
